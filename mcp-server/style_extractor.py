"""
style_extractor.py
既存のPPTXファイルからデザインスタイル情報を抽出するモジュール
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


def _rgb_to_hex(rgb: RGBColor | None) -> str | None:
    """RGBColorをHEX文字列に変換"""
    if rgb is None:
        return None
    return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"


def _safe_pt(emu_val: Any) -> float | None:
    """EMU値をポイントに変換"""
    try:
        return round(Pt(0).__class__(emu_val).pt, 1)  # type: ignore[attr-defined]
    except Exception:
        try:
            return round(emu_val / 12700, 1)
        except Exception:
            return None


def extract_style(pptx_path: str) -> dict[str, Any]:
    """
    PPTXファイルからスタイル情報を抽出してdictで返す

    Args:
        pptx_path: PPTXファイルのパス

    Returns:
        {
          "slide_width_emu": int,
          "slide_height_emu": int,
          "theme_colors": {...},
          "fonts": {"title": ..., "body": ...},
          "backgrounds": [...],
          "slide_layouts": [...],
          "sample_text_styles": [...]
        }
    """
    path = Path(pptx_path)
    if not path.exists():
        msg = f"PPTXファイルが見つかりません: {pptx_path}"
        raise FileNotFoundError(msg)

    prs = Presentation(str(path))
    style: dict[str, Any] = {}

    # スライドサイズ
    style["slide_width_emu"] = prs.slide_width
    style["slide_height_emu"] = prs.slide_height
    style["slide_width_pt"] = round(prs.slide_width / 12700, 1)
    style["slide_height_pt"] = round(prs.slide_height / 12700, 1)

    # テーマカラー抽出
    theme_colors = {}
    try:
        theme_element = prs.slide_master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )
        if theme_element is None:
            # slide_masterのXMLから直接取得
            import re

            xml_str = prs.slide_master.element.xml
            color_matches = re.findall(r'<a:(\w+)\s+val="([0-9A-Fa-f]{6})"', xml_str)
            for name, val in color_matches[:10]:
                theme_colors[name] = f"#{val.upper()}"
    except Exception:
        pass

    # カラーマップから主要色を取得
    try:
        slide_master = prs.slide_master
        # 背景色
        bg = slide_master.background
        if bg and bg.fill:
            fill = bg.fill
            if hasattr(fill, "fore_color") and fill.fore_color:
                try:
                    theme_colors["background"] = _rgb_to_hex(fill.fore_color.rgb)
                except Exception:
                    pass
    except Exception:
        pass

    style["theme_colors"] = theme_colors

    # フォント情報抽出（スライドマスターから）
    fonts = {"title": None, "body": None, "detected": []}
    try:
        for shape in prs.slide_master.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        font_info = {
                            "name": font.name,
                            "size_pt": _safe_pt(font.size) if font.size else None,
                            "bold": font.bold,
                            "color_hex": None,
                        }
                        try:
                            font_info["color_hex"] = _rgb_to_hex(font.color.rgb)
                        except Exception:
                            pass
                        if font_info["name"] and font_info not in fonts["detected"]:
                            fonts["detected"].append(font_info)
    except Exception:
        pass

    # slide_master のタイトル/本文フォント名を取得
    try:
        tx_styles = prs.slide_master.element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}txStyles"
        )
        if tx_styles is not None:
            import re

            ts_xml = tx_styles.xml if hasattr(tx_styles, "xml") else ""
            latin_matches = re.findall(r'<a:latin[^>]+typeface="([^"]+)"', ts_xml)
            if latin_matches:
                fonts["title"] = latin_matches[0]
                if len(latin_matches) > 1:
                    fonts["body"] = latin_matches[1]
    except Exception:
        pass

    style["fonts"] = fonts

    # スライドレイアウト情報
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders),
            "placeholders": [],
        }
        for ph in layout.placeholders:
            ph_info = {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
            }
            layout_info["placeholders"].append(ph_info)
        layouts.append(layout_info)
    style["slide_layouts"] = layouts

    # 実際のスライドからテキストスタイルサンプルを抽出（最大3枚）
    sample_styles = []
    for slide_idx, slide in enumerate(list(prs.slides)[:3]):
        slide_sample = {"slide_index": slide_idx, "shapes": []}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_info = {
                "shape_name": shape.name,
                "paragraphs": [],
            }
            for para in shape.text_frame.paragraphs[:3]:
                para_info = {
                    "text_preview": para.text[:50],
                    "alignment": str(para.alignment) if para.alignment else None,
                    "runs": [],
                }
                for run in para.runs[:2]:
                    font = run.font
                    run_info = {
                        "font_name": font.name,
                        "font_size_pt": _safe_pt(font.size) if font.size else None,
                        "bold": font.bold,
                        "color_hex": None,
                    }
                    try:
                        run_info["color_hex"] = _rgb_to_hex(font.color.rgb)
                    except Exception:
                        pass
                    para_info["runs"].append(run_info)
                shape_info["paragraphs"].append(para_info)
            slide_sample["shapes"].append(shape_info)
        sample_styles.append(slide_sample)
    style["sample_text_styles"] = sample_styles

    # 背景色サンプル（最大5枚）
    backgrounds = []
    for slide in list(prs.slides)[:5]:
        bg_info: dict[str, Any] = {"fill_type": None, "color_hex": None}
        try:
            fill = slide.background.fill
            bg_info["fill_type"] = str(fill.type)
            bg_info["color_hex"] = _rgb_to_hex(fill.fore_color.rgb)
        except Exception:
            pass
        backgrounds.append(bg_info)
    style["backgrounds"] = backgrounds

    return style


def extract_style_to_json(pptx_path: str, output_path: str | None = None) -> str:
    """
    PPTXからスタイル情報を抽出してJSONとして返す（オプションでファイル保存）

    Args:
        pptx_path: PPTXファイルのパス
        output_path: JSON出力先（Noneの場合はファイル保存しない）

    Returns:
        JSON文字列
    """
    style = extract_style(pptx_path)
    json_str = json.dumps(style, ensure_ascii=False, indent=2)

    if output_path:
        Path(output_path).write_text(json_str, encoding="utf-8")

    return json_str


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("使用法: python style_extractor.py <pptx_path> [output_json_path]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    result = extract_style_to_json(pptx_path, output_path)
    print(result)
