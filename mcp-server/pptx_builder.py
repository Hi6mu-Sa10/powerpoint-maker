"""
pptx_builder.py
スライド構成データ（dict）を受け取り、python-pptxでPPTXファイルを生成するモジュール
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────
# デフォルトカラーパレット（スタイル未指定時）
# ─────────────────────────────────────────────
DEFAULT_PALETTE = {
    "bg_title": RGBColor(0x1F, 0x35, 0x64),      # ダークネイビー
    "bg_section": RGBColor(0x2E, 0x74, 0xB5),     # ミッドブルー
    "bg_content": RGBColor(0xFF, 0xFF, 0xFF),      # ホワイト
    "accent": RGBColor(0x00, 0xB0, 0xF0),          # ライトブルー
    "text_light": RGBColor(0xFF, 0xFF, 0xFF),      # ホワイトテキスト
    "text_dark": RGBColor(0x1F, 0x35, 0x64),       # ダークネイビーテキスト
    "text_sub": RGBColor(0x40, 0x40, 0x40),        # ダークグレー
    "text_bullet": RGBColor(0x26, 0x26, 0x26),     # ほぼ黒
    "divider": RGBColor(0x00, 0xB0, 0xF0),         # アクセントブルー
}

DEFAULT_FONTS = {
    "title": "メイリオ",
    "body": "メイリオ",
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """HEX文字列をRGBColorに変換"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _build_palette(style: dict[str, Any] | None) -> dict[str, RGBColor]:
    """スタイル情報からカラーパレットを構築"""
    palette = dict(DEFAULT_PALETTE)
    if not style:
        return palette

    theme_colors = style.get("theme_colors", {})
    # テーマカラーがあれば上書き
    if theme_colors.get("dk1"):
        palette["text_dark"] = _hex_to_rgb(theme_colors["dk1"])
    if theme_colors.get("lt1"):
        palette["text_light"] = _hex_to_rgb(theme_colors["lt1"])
    if theme_colors.get("accent1"):
        palette["accent"] = _hex_to_rgb(theme_colors["accent1"])
        palette["divider"] = _hex_to_rgb(theme_colors["accent1"])

    return palette


def _build_fonts(style: dict[str, Any] | None) -> dict[str, str]:
    """スタイル情報からフォント設定を構築"""
    fonts = dict(DEFAULT_FONTS)
    if not style:
        return fonts
    f = style.get("fonts", {})
    if f.get("title"):
        fonts["title"] = f["title"]
    if f.get("body"):
        fonts["body"] = f["body"]
    return fonts


def _set_bg_color(slide: Any, color: RGBColor) -> None:
    """スライド背景色を設定"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> Any:
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_rect(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    line_color: RGBColor | None = None,
) -> Any:
    """矩形を追加"""
    shape = slide.shapes.add_shape(
        1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ─────────────────────────────────────────────
# スライドタイプ別レンダラー
# ─────────────────────────────────────────────

def _render_title_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """タイトルスライドを生成"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    _set_bg_color(slide, palette["bg_title"])

    # アクセントバー（左端の縦線）
    _add_rect(slide, 0.5, 1.8, 0.08, 3.2, palette["accent"])

    # タイトル
    _add_text_box(
        slide,
        data.get("title", "タイトルをここに入力"),
        0.75, 1.7, 8.5, 1.6,
        fonts["title"], 36,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )

    # サブタイトル
    _add_text_box(
        slide,
        data.get("subtitle", "サブタイトル・概要をここに入力"),
        0.75, 3.5, 8.5, 1.0,
        fonts["body"], 20,
        palette["accent"], bold=False,
        align=PP_ALIGN.LEFT,
    )

    # 区切り線
    _add_rect(slide, 0.75, 4.7, 8.5, 0.04, palette["accent"])

    # 日付・発表者
    _add_text_box(
        slide,
        data.get("meta", "20XX年XX月　発表者氏名　所属部署"),
        0.75, 4.85, 8.5, 0.5,
        fonts["body"], 14,
        RGBColor(0xCC, 0xCC, 0xCC), bold=False,
        align=PP_ALIGN.LEFT,
    )


def _render_agenda_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """目次スライドを生成"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _set_bg_color(slide, palette["bg_content"])

    # ヘッダーバー
    _add_rect(slide, 0, 0, 10.0, 1.3, palette["bg_title"])

    # スライドタイトル
    _add_text_box(
        slide,
        data.get("title", "目次"),
        0.4, 0.15, 9.0, 1.0,
        fonts["title"], 28,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )

    # アクセントライン
    _add_rect(slide, 0, 1.3, 10.0, 0.06, palette["accent"])

    items = data.get("items", [])
    start_y = 1.65
    item_h = 0.65

    for i, item in enumerate(items):
        y = start_y + i * item_h
        num_str = str(item.get("number", i + 1)).zfill(2)

        # 番号バッジ
        _add_rect(slide, 0.5, y, 0.5, 0.45, palette["bg_section"])
        _add_text_box(
            slide, num_str,
            0.5, y, 0.5, 0.45,
            fonts["title"], 14,
            palette["text_light"], bold=True,
            align=PP_ALIGN.CENTER,
        )

        # 項目テキスト
        _add_text_box(
            slide,
            item.get("text", f"アジェンダ {i+1}"),
            1.2, y, 8.0, 0.45,
            fonts["body"], 18,
            palette["text_dark"], bold=False,
            align=PP_ALIGN.LEFT,
        )


def _render_section_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """セクション区切りスライドを生成"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _set_bg_color(slide, palette["bg_section"])

    # 装飾的な矩形
    _add_rect(slide, 0, 5.5, 10.0, 2.0, palette["bg_title"])
    _add_rect(slide, 0, 3.3, 0.15, 2.3, palette["accent"])

    # セクション番号
    section_num = data.get("section_number", "")
    if section_num:
        _add_text_box(
            slide, f"SECTION  {section_num}",
            0.5, 2.2, 9.0, 0.6,
            fonts["body"], 14,
            RGBColor(0xCC, 0xEE, 0xFF), bold=False,
            align=PP_ALIGN.LEFT,
        )

    # セクションタイトル
    _add_text_box(
        slide,
        data.get("title", "セクションタイトル"),
        0.5, 2.9, 9.0, 1.8,
        fonts["title"], 36,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )


def _render_content_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """コンテンツ（箇条書き）スライドを生成"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _set_bg_color(slide, palette["bg_content"])

    # ヘッダーバー
    _add_rect(slide, 0, 0, 10.0, 1.2, palette["bg_title"])
    _add_rect(slide, 0, 1.2, 10.0, 0.06, palette["accent"])

    # スライドタイトル
    _add_text_box(
        slide,
        data.get("title", "スライドタイトル"),
        0.4, 0.1, 9.0, 1.0,
        fonts["title"], 26,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )

    bullets = data.get("bullets", [])
    start_y = 1.45
    spacing = 0.72

    for i, bullet in enumerate(bullets):
        y = start_y + i * spacing
        # バレットマーク
        _add_rect(slide, 0.4, y + 0.15, 0.12, 0.12, palette["accent"])

        # バレットテキスト
        if isinstance(bullet, dict):
            text = bullet.get("text", "")
            sub = bullet.get("sub", "")
        else:
            text = str(bullet)
            sub = ""

        _add_text_box(
            slide, text,
            0.65, y, 9.0, 0.42,
            fonts["body"], 16,
            palette["text_bullet"], bold=False,
            align=PP_ALIGN.LEFT,
        )
        if sub:
            _add_text_box(
                slide, f"  → {sub}",
                0.85, y + 0.38, 8.8, 0.32,
                fonts["body"], 12,
                palette["text_sub"], bold=False,
                align=PP_ALIGN.LEFT,
            )

    # ノートプレースホルダー
    note = data.get("note", "")
    if note:
        _add_rect(slide, 0.4, 6.9, 9.2, 0.04, palette["divider"])
        _add_text_box(
            slide, f"📝 {note}",
            0.4, 6.95, 9.2, 0.4,
            fonts["body"], 10,
            palette["text_sub"], bold=False,
            align=PP_ALIGN.LEFT,
        )


def _render_two_column_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """2カラムスライドを生成"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _set_bg_color(slide, palette["bg_content"])

    # ヘッダーバー
    _add_rect(slide, 0, 0, 10.0, 1.2, palette["bg_title"])
    _add_rect(slide, 0, 1.2, 10.0, 0.06, palette["accent"])

    # スライドタイトル
    _add_text_box(
        slide,
        data.get("title", "スライドタイトル"),
        0.4, 0.1, 9.0, 1.0,
        fonts["title"], 26,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )

    left_col = data.get("left", {})
    right_col = data.get("right", {})

    for col_data, x_offset in [(left_col, 0.4), (right_col, 5.2)]:
        # カラムヘッダー
        col_title = col_data.get("title", "")
        if col_title:
            _add_rect(slide, x_offset, 1.4, 4.5, 0.45, palette["bg_section"])
            _add_text_box(
                slide, col_title,
                x_offset, 1.4, 4.5, 0.45,
                fonts["title"], 14,
                palette["text_light"], bold=True,
                align=PP_ALIGN.CENTER,
            )

        items = col_data.get("items", [])
        for j, item in enumerate(items):
            y = 2.05 + j * 0.65
            _add_rect(slide, x_offset + 0.05, y + 0.15, 0.10, 0.10, palette["accent"])
            _add_text_box(
                slide, str(item),
                x_offset + 0.25, y, 4.2, 0.55,
                fonts["body"], 14,
                palette["text_bullet"], bold=False,
                align=PP_ALIGN.LEFT,
            )

    # 中央区切り線
    _add_rect(slide, 4.9, 1.4, 0.04, 5.7, palette["divider"])


def _render_timeline_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """タイムライン・ロードマップスライドを生成"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _set_bg_color(slide, palette["bg_content"])

    # ヘッダーバー
    _add_rect(slide, 0, 0, 10.0, 1.2, palette["bg_title"])
    _add_rect(slide, 0, 1.2, 10.0, 0.06, palette["accent"])

    # スライドタイトル
    _add_text_box(
        slide,
        data.get("title", "ロードマップ"),
        0.4, 0.1, 9.0, 1.0,
        fonts["title"], 26,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )

    milestones = data.get("milestones", [])
    if not milestones:
        return

    col_w = 9.0 / max(len(milestones), 1)

    # 中央ライン
    _add_rect(slide, 0.4, 3.5, 9.0, 0.06, palette["bg_section"])

    for i, ms in enumerate(milestones):
        x = 0.4 + i * col_w
        center_x = x + col_w / 2 - 0.3

        # フェーズラベル
        _add_rect(slide, center_x, 1.45, 0.6, 1.9, palette["bg_section"])
        _add_text_box(
            slide, ms.get("phase", f"P{i+1}"),
            center_x, 1.5, 0.6, 1.8,
            fonts["title"], 11,
            palette["text_light"], bold=True,
            align=PP_ALIGN.CENTER,
        )

        # 丸マーカー
        _add_rect(slide, center_x + 0.07, 3.35, 0.46, 0.42, palette["accent"])

        # 期間ラベル
        _add_text_box(
            slide, ms.get("period", ""),
            x, 3.85, col_w, 0.35,
            fonts["body"], 11,
            palette["text_dark"], bold=True,
            align=PP_ALIGN.CENTER,
        )

        # マイルストーン内容
        _add_text_box(
            slide, ms.get("description", ""),
            x, 4.3, col_w, 1.8,
            fonts["body"], 11,
            palette["text_sub"], bold=False,
            align=PP_ALIGN.CENTER,
        )


def _render_closing_slide(
    prs: Presentation,
    data: dict[str, Any],
    palette: dict[str, RGBColor],
    fonts: dict[str, str],
) -> None:
    """クロージングスライドを生成"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _set_bg_color(slide, palette["bg_title"])

    _add_rect(slide, 0.5, 2.8, 0.1, 2.0, palette["accent"])

    _add_text_box(
        slide,
        data.get("title", "Thank You"),
        0.75, 2.7, 8.5, 1.4,
        fonts["title"], 44,
        palette["text_light"], bold=True,
        align=PP_ALIGN.LEFT,
    )

    _add_rect(slide, 0.75, 4.25, 8.5, 0.05, palette["accent"])

    _add_text_box(
        slide,
        data.get("message", "ご質問・ご意見はお気軽にどうぞ"),
        0.75, 4.4, 8.5, 0.6,
        fonts["body"], 18,
        palette["accent"], bold=False,
        align=PP_ALIGN.LEFT,
    )

    contact = data.get("contact", "")
    if contact:
        _add_text_box(
            slide, contact,
            0.75, 5.2, 8.5, 0.6,
            fonts["body"], 14,
            RGBColor(0xCC, 0xCC, 0xCC), bold=False,
            align=PP_ALIGN.LEFT,
        )


# ─────────────────────────────────────────────
# メインビルダー
# ─────────────────────────────────────────────

SLIDE_RENDERERS = {
    "title": _render_title_slide,
    "agenda": _render_agenda_slide,
    "section": _render_section_slide,
    "content": _render_content_slide,
    "two_column": _render_two_column_slide,
    "timeline": _render_timeline_slide,
    "closing": _render_closing_slide,
}


def build_pptx(
    slides_data: list[dict[str, Any]],
    output_path: str,
    style: dict[str, Any] | None = None,
) -> str:
    """
    スライドデータからPPTXを生成

    Args:
        slides_data: スライドのリスト。各要素は {"type": "...", ...} の形式
        output_path: 出力先PPTXファイルパス
        style: style_extractor.py で抽出したスタイル辞書（省略可）

    Returns:
        生成されたファイルの絶対パス
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    palette = _build_palette(style)
    fonts = _build_fonts(style)

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        renderer = SLIDE_RENDERERS.get(slide_type, _render_content_slide)
        renderer(prs, slide_data, palette, fonts)

    # 出力ディレクトリ作成
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs.save(str(out))
    return str(out.resolve())


if __name__ == "__main__":
    # 動作テスト
    sample_slides: list[dict[str, Any]] = [
        {
            "type": "title",
            "title": "AI活用による業務効率化提案",
            "subtitle": "生成AIを用いた次世代ワークフロー構築",
            "meta": "2025年5月　システム企画部　山田太郎",
        },
        {
            "type": "agenda",
            "title": "目次",
            "items": [
                {"number": 1, "text": "現状の課題"},
                {"number": 2, "text": "提案するソリューション"},
                {"number": 3, "text": "期待される効果"},
                {"number": 4, "text": "実装スケジュール"},
                {"number": 5, "text": "まとめ"},
            ],
        },
        {
            "type": "section",
            "section_number": "01",
            "title": "現状の課題",
        },
        {
            "type": "content",
            "title": "現状の課題",
            "bullets": [
                {"text": "手作業によるデータ入力に多大な時間を消費", "sub": "1日あたり約3時間の工数"},
                {"text": "情報の分散により意思決定が遅延", "sub": "部署間の連携ロスが発生"},
                {"text": "ナレッジの属人化リスク", "sub": "退職・異動時の業務継続性に課題"},
                {"text": "レポート作成の煩雑さ", "sub": "週次レポートに平均2時間"},
            ],
            "note": "出典: 社内業務時間調査（2024年Q4）",
        },
        {
            "type": "two_column",
            "title": "AI導入前後の比較",
            "left": {
                "title": "導入前（現状）",
                "items": ["手動でのデータ集計", "メール・Excel中心の管理", "レポートは手作業", "ナレッジが個人管理"],
            },
            "right": {
                "title": "導入後（目標）",
                "items": ["自動データ収集・集計", "統合ダッシュボード管理", "レポート自動生成", "AI活用ナレッジ共有"],
            },
        },
        {
            "type": "timeline",
            "title": "実装ロードマップ",
            "milestones": [
                {"phase": "Phase 1", "period": "2025年 Q3", "description": "要件定義・\nPoC実施"},
                {"phase": "Phase 2", "period": "2025年 Q4", "description": "パイロット\n導入"},
                {"phase": "Phase 3", "period": "2026年 Q1", "description": "全社展開\n開始"},
                {"phase": "Phase 4", "period": "2026年 Q2", "description": "効果測定・\n改善"},
            ],
        },
        {
            "type": "closing",
            "title": "Thank You",
            "message": "ご質問・ご意見はお気軽にどうぞ",
            "contact": "システム企画部　yamada@example.com　|　内線: 1234",
        },
    ]

    output = build_pptx(sample_slides, "output/sample.pptx")
    print(f"✅ PPTXを生成しました: {output}")
